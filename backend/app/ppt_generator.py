from pptx import Presentation
from pptx.util import Inches, Pt
import re
from app.image_service import search_images, get_fallback_images
from app.ai_service import generate_ppt_outline, is_api_available


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


TITLE_SLIDE_LAYOUT = 0
TITLE_CONTENT_LAYOUT = 1
BLANK_LAYOUT = 6


def generate_ppt(text: str, output_path: str, use_ai: bool = True):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    if use_ai and is_api_available():
        generate_ppt_with_ai(prs, text)
    else:
        generate_ppt_fallback(prs, text)

    prs.save(output_path)


def generate_ppt_with_ai(prs: Presentation, text: str):
    outline = generate_ppt_outline(text)

    if outline and len(outline) > 0:
        first_slide = outline[0]
        add_title_slide(prs, first_slide.get("title", "演示文稿"))

        for slide_data in outline[1:]:
            title = slide_data.get("title", "")
            content = slide_data.get("content", [])
            if isinstance(content, str):
                content = [c.strip() for c in content.split('\n') if c.strip()]
            add_content_slide(prs, title, content[:5])
    else:
        generate_ppt_fallback(prs, text)


def generate_ppt_fallback(prs: Presentation, text: str):
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    title = lines[0] if lines else "演示文稿"
    add_title_slide(prs, title)

    content_lines = lines[1:] if len(lines) > 1 else []
    sections = split_into_sections(content_lines)

    for section_title, section_content in sections:
        add_content_slide(prs, section_title, section_content)


def add_title_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE_LAYOUT])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True


def add_content_slide(prs: Presentation, title: str, content: list):
    slide = prs.slides.add_slide(prs.slide_layouts[TITLE_CONTENT_LAYOUT])

    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True

    content_shape = slide.placeholders[1]
    tf = content_shape.text_frame
    tf.clear()

    for i, line in enumerate(content):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {line}"
        p.font.size = Pt(20)
        p.level = 0


def split_into_sections(lines: list) -> list:
    sections = []
    current_section = None
    current_content = []

    for line in lines:
        if len(line) < 50 and not line.startswith('•') and not line.startswith('-'):
            if current_section:
                sections.append((current_section, current_content))
            current_section = line
            current_content = []
        else:
            clean_line = re.sub(r'^[•\-\s]+', '', line).strip()
            if clean_line:
                current_content.append(clean_line)

    if current_section and current_content:
        sections.append((current_section, current_content))
    elif current_content:
        sections.append(("内容", current_content))

    if not sections:
        sections = [("概述", lines[:5])]

    return sections[:10]